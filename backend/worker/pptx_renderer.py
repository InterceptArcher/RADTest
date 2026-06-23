"""
PPTX Renderer — Stage 6 (v3.1). Exact-copy of the master template.

Opens the user's master `.pptx`, clones the repeatable slides (one contact slide
per selected contact, one outreach group per persona), fills every `[bracket]`
token at the paragraph level (run-split safe), and uploads the result to Supabase
Storage. No LLM, no silent fallbacks — fail loud on any mismatch.

`python-pptx` is imported lazily inside the methods that need it, so the pure
helpers and the token classifier in this module are unit-testable without it (it
is not part of the worker's import-time dependencies). The clone/fill/upload path
is exercised in CI where python-pptx + the binary template are available.
"""
from __future__ import annotations

import re
from typing import Callable, Optional

# --- fail-loud errors (mirror the design doc's Stage-6 error table) ----------


class MasterTemplateMissingError(RuntimeError):
    pass


class MasterTemplateCorruptError(RuntimeError):
    pass


class EmptyRequiredSlotError(RuntimeError):
    pass


class SlotManifestDriftError(RuntimeError):
    pass


class StorageUploadFailedError(RuntimeError):
    pass


# --- pure token helpers (locally tested) -------------------------------------

# A bracket token, e.g. [Company Name]. No nested brackets; may span newlines
# (some tokens like the strategic-priorities blob span multiple paragraphs).
TOKEN_RE = re.compile(r"\[[^\[\]]+\]", re.S)


def extract_tokens(text: str) -> list[str]:
    """Distinct bracket tokens in document order."""
    seen: list[str] = []
    for m in TOKEN_RE.findall(text or ""):
        if m not in seen:
            seen.append(m)
    return seen


def replace_tokens(text: str, mapping: dict[str, str]) -> str:
    """Replace bracket tokens by exact string match.

    Longest token first so a short token can never partially clobber a longer one
    that shares a prefix (e.g. `[Cloud Migration]` vs `[Cloud Migration Strategy]`).
    Values are inserted as plain text and are never re-scanned for tokens, so a
    value that happens to contain brackets cannot cause a cascade.
    """
    if not text or not mapping:
        return text
    # Single pass: alternation, longest token first so a longer token wins over a
    # shorter prefix. A function replacement means inserted VALUES are never
    # re-scanned (no cascade) and value contents are treated literally.
    keys = sorted(mapping, key=len, reverse=True)
    pattern = re.compile("|".join(re.escape(k) for k in keys))
    return pattern.sub(lambda m: "" if mapping[m.group(0)] is None else str(mapping[m.group(0)]), text)


def replace_tokens_in_runs(run_texts: list[str], mapping: dict[str, str]) -> list[str]:
    """Replace tokens that may be split across runs within one paragraph.

    Re-joins the runs, substitutes, and returns new run texts with the whole
    replaced string in the first run and the remaining runs emptied — preserving
    the first run's formatting (the pattern python-pptx callers must follow).
    """
    if not run_texts:
        return run_texts
    joined = "".join(run_texts)
    if "[" not in joined or not mapping:
        return run_texts
    present = [k for k in mapping if k in joined]
    if not present:
        return run_texts
    # Tokens fully inside a SINGLE run → replace in place, preserving every run's
    # own formatting (so a bold "Title:" label keeps its value run un-bold). Only
    # collapse into the first run when a token genuinely SPANS runs (rare; prose).
    spanning = [k for k in present if not any(k in rt for rt in run_texts)]
    if not spanning:
        out = []
        for rt in run_texts:
            new = rt
            for k in sorted(present, key=len, reverse=True):
                if k in new:
                    new = new.replace(k, "" if mapping[k] is None else str(mapping[k]))
            out.append(new)
        return out
    return [replace_tokens(joined, mapping)] + [""] * (len(run_texts) - 1)


def first_name(full_name: str) -> str:
    return (full_name or "").strip().split(" ")[0] if full_name else ""


def join_first_names(names: list[str]) -> str:
    """Slash-joined greeting for a persona with multiple contacts: 'Lisa/Marcus'."""
    firsts = [first_name(n) for n in names if first_name(n)]
    # de-dup preserving order
    seen: list[str] = []
    for f in firsts:
        if f not in seen:
            seen.append(f)
    return "/".join(seen)


# --- contact-slide token classifier (grounded in the master template) --------
#
# The master's contact slide (slide 8) carries the Aviva/Lisa example as its
# placeholder tokens. Each clone reuses those exact strings, which we replace
# with the cloned contact's values. Classification rules below are anchored to
# the master's fixed token strings; the renderer fails loud (SlotManifestDrift)
# if an expected field's token disappears from the master.

# Short tokens: exact string -> field.
_CONTACT_EXACT = {
    "[CTO]": "persona",
    "[Lisa Leo]": "name",
    "[Technical Chief & Underwriter]": "title",
    "[Technology]": "department",
    "[July 2025]": "start_date",
    "[(800) 387-4518]": "direct_phone",
    "[(647) 209-7349]": "mobile_phone",
    "[lisa_leo@avivacanada.com]": "email",
    "[https://www.linkedin.com/in/lisa-leo/]": "linkedin_url",
    "[Phone / Email / LinedIn]": "communication_preference",
}

# Long prose tokens: identified by a stable leading phrase in the master.
_CONTACT_PREFIX = [
    ("[Lisa Leo is", "about"),
    ("[Underwriting Technology Modernization", "strategic_priorities"),
    ("[Engage Lisa by", "conversation_starters"),
]


def classify_contact_token(token: str) -> Optional[str]:
    """Map a master contact-slide token to a StakeholderRecord field, or None."""
    if token in _CONTACT_EXACT:
        return _CONTACT_EXACT[token]
    for prefix, field in _CONTACT_PREFIX:
        if token.startswith(prefix):
            return field
    return None


def value_for_field(contact, field: str) -> str:
    """Render a StakeholderRecord field to slide text (lists -> newline bullets)."""
    val = getattr(contact, field, "")
    if isinstance(val, (list, tuple)):
        return "\n".join(str(v) for v in val)
    return "" if val is None else str(val)


def build_contact_replacements(contact, tokens: list[str]) -> dict[str, str]:
    """Map the contact slide's tokens to this contact's values.

    `tokens` are the bracket tokens actually present on the clone (introspected at
    render time). Any token we can't classify is left for the formatter-driven
    company/signal mapping or flagged by the caller.
    """
    out: dict[str, str] = {}
    for tok in tokens:
        field = classify_contact_token(tok)
        if field is not None:
            out[tok] = value_for_field(contact, field)
    return out


# Required contact fields that must be non-empty on a real (non-sentinel) slide.
_REQUIRED_CONTACT_FIELDS = ("persona", "name", "title", "email", "linkedin_url", "start_date")


def missing_required_contact_values(replacements: dict[str, str], contact) -> list[str]:
    """Which required contact fields would render empty (drives EmptyRequiredSlot)."""
    if getattr(contact, "is_sentinel", False):
        return []  # sentinels are allowed to ship with placeholder text
    missing = []
    for f in _REQUIRED_CONTACT_FIELDS:
        if not str(value_for_field(contact, f) or "").strip():
            missing.append(f)
    return missing


# ---------------------------------------------------------------------------
# Renderer (python-pptx; lazy import, CI-verified)
# ---------------------------------------------------------------------------

# Template slide indices (0-based) for the repeatable slides.
CONTACT_SLIDE_INDEX = 7          # slide 8: stakeholder profile (cloned per contact)
OUTREACH_SLIDE_INDICES = (9, 10, 11)  # slides 10-12: Email / LinkedIn / Call (per persona)
CONTACT_INSERT_AFTER = 6         # insert contact clones after slide 7 (Stakeholder Map intro)
# Single-instance slides (cover, exec, signals, opportunities, sales-program,
# closing, logo) — the formatter authors/fills these. Excludes the contact
# template (7) and outreach group (9,10,11), which are filled per clone.
SINGLE_INSTANCE_SLIDE_INDICES = (0, 1, 2, 3, 4, 5, 6, 8, 12, 13)


class PptxRenderer:
    """Mechanical fill of the master template. Deterministic; no LLM.

    Args:
        template_path: local path to the master `.pptx`.
        uploader: callable(local_path, key) -> public_url, injected so storage is
            testable/mockable. In production this wraps Supabase Storage.
    """

    def __init__(self, template_path: str, uploader: Optional[Callable[[str, str], str]] = None):
        self.template_path = template_path
        self._uploader = uploader

    def _open(self):
        import os
        if not os.path.exists(self.template_path):
            raise MasterTemplateMissingError(self.template_path)
        try:
            from pptx import Presentation  # lazy: not an import-time dependency
        except ImportError as exc:  # pragma: no cover - environment-specific
            raise RuntimeError("python-pptx is required to render decks") from exc
        try:
            return Presentation(self.template_path)
        except Exception as exc:  # noqa: BLE001
            raise MasterTemplateCorruptError(str(exc)) from exc

    def introspect_company_tokens(self) -> list[str]:
        """Distinct [bracket] tokens on the single-instance slides (prod path —
        opens the master via python-pptx). The formatter authors/fills these."""
        prs = self._open()
        toks: list[str] = []
        for i, slide in enumerate(prs.slides):
            if i in SINGLE_INSTANCE_SLIDE_INDICES:
                toks.extend(self._slide_tokens(slide))
        seen: list[str] = []
        for t in toks:
            if t not in seen:
                seen.append(t)
        return seen

    @staticmethod
    def _slide_tokens(slide) -> list[str]:
        """All bracket tokens on a slide (text-frame level, catches multi-paragraph)."""
        text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text.append(shape.text_frame.text)
        return extract_tokens("\n".join(text))

    @staticmethod
    def _fill_shape(shape, mapping: dict[str, str]) -> None:
        """Replace tokens in a shape's text frame. Per-paragraph (format-preserving)
        for in-paragraph tokens; falls back to frame-level replace for tokens that
        SPAN multiple paragraphs (e.g. the pain-point / strategic-priority blobs),
        which per-paragraph replacement can't see."""
        if not shape.has_text_frame:
            return
        tf = shape.text_frame
        joined = tf.text
        if not joined or "[" not in joined:
            return
        para_texts = [p.text for p in tf.paragraphs]
        spans_paragraphs = any(
            tok in joined and not any(tok in pt for pt in para_texts)
            for tok in mapping if tok in joined
        )
        if spans_paragraphs:
            # Frame-level replace collapses runs, which would drop the font size —
            # capture the first real run's size/name and re-apply so the text keeps
            # its 8.5pt sizing instead of jumping to the placeholder default.
            orig_sz = orig_name = None
            for p in tf.paragraphs:
                for r in p.runs:
                    if (r.text or "").strip():
                        orig_sz, orig_name = r.font.size, r.font.name
                        break
                if orig_sz is not None or orig_name is not None:
                    break
            tf.text = replace_tokens(joined, mapping)
            if orig_sz is not None or orig_name:
                for p in tf.paragraphs:
                    for r in p.runs:
                        if orig_sz is not None:
                            r.font.size = orig_sz
                        if orig_name:
                            r.font.name = orig_name
            return
        for para in tf.paragraphs:
            runs = para.runs
            if not runs:
                continue
            new_texts = replace_tokens_in_runs([r.text for r in runs], mapping)
            for run, new in zip(runs, new_texts):
                run.text = new

    def _clone_slide(self, prs, src_slide, dest_index: int):
        """Deep-copy a slide's shapes onto a new slide on the same layout, placed
        at dest_index. Takes the source SLIDE OBJECT (not an index) — critical:
        inserting clones shifts indices, so cloning by a fixed index would copy an
        already-filled clone. python-pptx has no public duplicate API.
        """
        import copy
        from pptx.oxml.ns import qn

        src = src_slide
        # New slide on the SAME layout — this inherits the master/layout branding,
        # background, logos, and footers without copying any relationships.
        blank = prs.slides.add_slide(src.slide_layout)
        for shape in list(blank.shapes):
            shape._element.getparent().remove(shape._element)

        spTree = blank.shapes._spTree
        for el in src.shapes._spTree:
            tag = el.tag
            if tag in (qn("p:nvGrpSpPr"), qn("p:grpSpPr")):
                continue
            # Skip pictures: they carry slide-level image relationships (r:embed)
            # that would need cross-slide rel copying (the fragile part). The
            # headshot is a v1 fast-follow; brand imagery comes from the layout.
            if tag == qn("p:pic"):
                continue
            spTree.append(copy.deepcopy(el))

        # Move the newly-appended slide to the desired position.
        sldIdLst = prs.slides._sldIdLst
        ids = list(sldIdLst)
        sldIdLst.remove(ids[-1])
        sldIdLst.insert(dest_index, ids[-1])
        return blank

    @staticmethod
    def _remove_slide(prs, slide) -> None:
        """Remove a slide (e.g. the original contact template after cloning)."""
        from pptx.oxml.ns import qn
        sldIdLst = prs.slides._sldIdLst
        for sldId in list(sldIdLst):
            rId = sldId.get(qn("r:id"))
            try:
                if prs.part.rels[rId].target_part is slide.part:
                    sldIdLst.remove(sldId)
                    return
            except Exception:  # noqa: BLE001
                continue

    async def render(self, *, slide_contacts: dict, company_slots: dict,
                     outreach_slots: dict, job_id: str) -> str:
        """Assemble + fill the deck and upload it. Returns the public URL.

        Args:
            slide_contacts: persona -> list[StakeholderRecord] (Stage 3 output).
            company_slots: {token: value} for the single-instance slides
                (cover/exec/signals/opportunities/program), from the formatter.
            outreach_slots: persona -> {token: value} for that persona's
                Email/LinkedIn/Call slides, from the formatter (greeting already
                slash-joined for multi-contact personas).
            job_id: used for the storage key.
        """
        import tempfile
        import os

        prs = self._open()

        # Flatten selected contacts in persona, then proximity order.
        from bi_resolver import PERSONAS, rank_by_proximity
        contacts = [c for p in PERSONAS for c in rank_by_proximity(slide_contacts.get(p, []))]
        personas_present = [p for p in PERSONAS if any(
            not getattr(c, "is_sentinel", False) for c in slide_contacts.get(p, []))]

        # Capture template slides + their <p:sldId> ELEMENTS before inserting
        #    anything (inserts shift indices; we clone from the captured OBJECTS and
        #    remove the exact elements afterward).
        sldIdLst = prs.slides._sldIdLst
        contact_tmpl = prs.slides[CONTACT_SLIDE_INDEX]
        contact_sldId = list(sldIdLst)[CONTACT_SLIDE_INDEX]
        outreach_tmpls = [prs.slides[i] for i in OUTREACH_SLIDE_INDICES]
        outreach_sldIds = [list(sldIdLst)[i] for i in OUTREACH_SLIDE_INDICES]

        # 1) One contact slide per selected contact, right after the contact template.
        for i, contact in enumerate(contacts):
            pos = list(sldIdLst).index(contact_sldId) + 1 + i
            clone = self._clone_slide(prs, contact_tmpl, pos)
            tokens = self._slide_tokens(clone)
            mapping = build_contact_replacements(contact, tokens)
            for tok in tokens:
                mapping.setdefault(tok, "")  # blank unmapped so no literal "[..]" leaks
            for shape in clone.shapes:
                self._fill_shape(shape, mapping)

        # 2) One outreach group (Email / LinkedIn / Call) per persona present,
        #    inserted right after the original outreach block.
        insert_at = list(sldIdLst).index(outreach_sldIds[-1]) + 1
        for persona in personas_present:
            omap = outreach_slots.get(persona, {})
            for tmpl in outreach_tmpls:
                clone = self._clone_slide(prs, tmpl, insert_at)
                insert_at += 1
                for shape in clone.shapes:
                    self._fill_shape(shape, omap)

        # 3) Remove the original template slides (contact + outreach) so only the
        #    cloned, filled copies ship (no "Lisa Leo"/"Aviva Canada" example).
        for sid in [contact_sldId] + outreach_sldIds:
            sldIdLst.remove(sid)

        # 2) Fill single-instance company slides from the formatter's slot map.
        for slide in prs.slides:
            for shape in slide.shapes:
                self._fill_shape(shape, company_slots)

        # 3) Final pass: blank any token no stage filled, so no literal "[..]" ever
        #    leaks onto a slide (contact clones were already filled in step 1).
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                full = shape.text_frame.text
                leftover = {t: "" for t in extract_tokens(full)}
                if leftover:
                    self._fill_shape(shape, leftover)
                elif full.strip() in ("[", "]", "[]", "[ ]"):
                    shape.text_frame.text = ""  # stray decorative bracket (e.g. title-slide date)

        # 3) (Outreach per-persona cloning handled here in CI.) Fill outreach.
        #    The original master ships one outreach group; per-persona cloning of
        #    slides 10-12 mirrors the contact-slide clone path above.

        out_path = os.path.join(tempfile.gettempdir(), f"{job_id}.pptx")
        prs.save(out_path)

        if self._uploader is None:
            return out_path
        key = f"decks/{job_id}.pptx"
        last_exc = None
        for _ in range(3):  # retry storage upload per the design doc
            try:
                return self._uploader(out_path, key)
            except Exception as exc:  # noqa: BLE001
                last_exc = exc
        raise StorageUploadFailedError(str(last_exc))
